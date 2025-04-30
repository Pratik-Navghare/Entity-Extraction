import pdfplumber
from pptx import Presentation
import pytesseract
from PIL import Image
import json
import os
import re

try:
    from langchain_ollama import OllamaLLM
except ImportError:
    print("Please install langchain-ollama package: pip install langchain-ollama")
    OllamaLLM = None

def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\\n"
    return text

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\\n"
    return text

def extract_text_from_image(image_path):
    image = Image.open(image_path)
    text = pytesseract.image_to_string(image)
    return text

def extract_entities_with_llama(text):
    if OllamaLLM is None:
        print("OllamaLLM is not available. Please install langchain-ollama.")
        return []
    prompt = f"Extract at least 20 entities are required from the following text and return as JSON array of key,type pair of entities and don't add any things outside the given text also is has to be JSON format:\\n{text}"
    try:
        llm = OllamaLLM(model="llama3.2")
        response = llm.invoke(prompt)
        print("Raw Ollama response:", response)
        try:
            entities = json.loads(response)
        except json.JSONDecodeError:
            print("Failed to parse JSON from Ollama response")
            json_match = re.search(r'```(?:json)?\s*(\[[\s\S]*?\])\s*```', response, re.DOTALL | re.IGNORECASE)
            if json_match:
                json_str = json_match.group(1)
                try:
                    entities = json.loads(json_str)
                except json.JSONDecodeError:
                    entities = []
            else:
                # Fallback: try to extract any JSON array in the response
                json_match = re.search(r'\[[\s\S]*\]', response)
                if json_match:
                    try:
                        entities = json.loads(json_match.group())
                    except json.JSONDecodeError:
                        entities = []
                else:
                    entities = []
            print("Failed to parse JSON from Ollama response")
            print("Failed to parse JSON from Ollama response")
        return entities
    except Exception as e:
        print(f"Error calling OllamaLLM via langchain-ollama: {e}")
        return []

def process_files(file_paths):
    combined_text = ""
    for file_path in file_paths:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            combined_text += extract_text_from_pdf(file_path) + "\\n"
        elif ext in [".png", ".jpeg", ".jpg"]:
            combined_text += extract_text_from_image(file_path) + "\\n"
        elif ext in [".ppt", ".pptx"]:
            combined_text += extract_text_from_pptx(file_path) + "\\n"
        else:
            print(f"Unsupported file type: {file_path}")
    entities = extract_entities_with_llama(combined_text)
    return entities
